import os
import shutil
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE
from services.data_context_engine import data_context_engine
from audit.top_risk_engine import get_top_5_risks, get_process_controls, DEFAULT_COL_MAP

# Path to templates
TEMPLATE_DIR = os.path.join(os.path.dirname(os.path.dirname(__file__)), 'templates')
OUTPUT_DIR = os.path.join(os.path.dirname(os.path.dirname(__file__)), 'uploads')

TEMPLATE_MAP = {
    "Internal Audit": "Internal Audit - Q3 2024.pptx",
    "SOX Compliance Review": "SOX Compliance Review - Q3 2024.pptx",
    "Process Review": "Process Review - Q3 2024.pptx",
    "Agile Ad-hoc Audit": "Agile Ad-hoc Audit - Q3 2024.pptx"
}

def load_template(audit_type: str):
    filename = TEMPLATE_MAP.get(audit_type, "Agile Ad-hoc Audit - Q3 2024.pptx")
    if audit_type == "Internal Audit":
        custom_fn = "Internal Audit - Q3 2 2024.pptx"
        if os.path.exists(os.path.join(TEMPLATE_DIR, custom_fn)):
            filename = custom_fn
            
    template_path = os.path.join(TEMPLATE_DIR, filename)
    
    if not os.path.exists(template_path):
        # Create a dummy template if it doesn't exist for testing purposes
        os.makedirs(TEMPLATE_DIR, exist_ok=True)
        prs = Presentation()
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = f"{audit_type} Kickoff"
        subtitle.text = "Audit Period: {{AUDIT_PERIOD}}\nManager: {{AUDIT_MANAGER}}\n\nScope:\n{{AUDIT_SCOPE}}\n\nTop Risks:\n{{TOP_RISKS}}\n\nProcesses:\n{{PROCESS_LIST}}\n\nKey Controls:\n{{KEY_CONTROLS}}"
        prs.save(template_path)
        
    return Presentation(template_path), filename

def generate_presentation(payload: dict) -> str:
    from pptx.util import Pt
    from pptx.dml.color import RGBColor

    audit_type = payload.get("audit_type", "Agile Ad-hoc Audit")
    audit_name = payload.get("audit_name", "Kickoff Presentation")
    start_date = payload.get("start_date", "N/A")
    end_date = payload.get("end_date", "N/A")
    
    auditors_raw = payload.get("auditors", [])
    if isinstance(auditors_raw, list):
        auditors = ", ".join(auditors_raw)
    else:
        auditors = str(auditors_raw)
        
    tenant_id = payload.get("tenant_id", "CJSJ")
    
    # Load template
    prs, filename = load_template(audit_type)
    
    # Extract real RCM data (needed for fallback)
    df = data_context_engine.get_dataframe(tenant_id)
    
    # Parse rows data
    payload_rows = payload.get("rows")
    processes_data = []
    
    if payload_rows is not None:
        # 1. Parse rows directly from frontend state
        for r in payload_rows:
            pname = r.get("process_name")
            if not pname:
                continue
            risk_area = r.get("risk_area") or "N/A"
            risk_score = r.get("risk_score") or 0
            man_hours = r.get("man_hours") or 0
            auditors_count = r.get("auditors") or 0
            
            # Extract controls
            nested_controls = []
            for ctrl in (r.get("controls") or []):
                if isinstance(ctrl, str):
                    desc = ctrl.strip()
                elif isinstance(ctrl, dict):
                    desc = (ctrl.get("control_description") or ctrl.get("ControlDescription") or "").strip()
                else:
                    desc = ""
                if desc:
                    nested_controls.append(desc)
                    
            processes_data.append({
                "process_name": pname,
                "risk_area": risk_area,
                "risk_score": risk_score,
                "man_hours": man_hours,
                "auditors_count": auditors_count,
                "controls": nested_controls
            })
    else:
        # 2. Fallback to reading from df
        if df is None or df.empty:
            raise ValueError("No RCM data uploaded. Please upload an RCM first.")
            
        proc_col = DEFAULT_COL_MAP.get("process", "process")
        if df is not None and proc_col in df.columns:
            unique_procs = df[proc_col].dropna().unique().tolist()
            for pname in unique_procs:
                p_df = df[df[proc_col] == pname]
                
                # Extract controls
                ctrl_col = DEFAULT_COL_MAP.get("control_description", "control_description")
                controls = p_df[ctrl_col].dropna().unique().tolist()[:10] if (ctrl_col and ctrl_col in p_df.columns) else []
                
                processes_data.append({
                    "process_name": pname,
                    "risk_area": "Derived Process Area",
                    "risk_score": 50,
                    "man_hours": 40,
                    "auditors_count": 1,
                    "controls": controls
                })
                
    # Modify Slide 1 (details)
    slide_1 = prs.slides[0]
    for shape in slide_1.shapes:
        if shape.has_text_frame:
            text_lower = shape.text_frame.text.lower()
            if any(t.lower() in text_lower for t in TEMPLATE_MAP.keys()) or "kickoff" in text_lower or text_lower == audit_type.lower():
                shape.text_frame.text = f"{audit_type} - {audit_name}"
                for paragraph in shape.text_frame.paragraphs:
                    paragraph.font.name = "Inter"
                    paragraph.font.bold = True
                    paragraph.font.size = Pt(24)
                    paragraph.font.color.rgb = RGBColor(15, 23, 42)
            elif "period" in text_lower or "audit period" in text_lower:
                shape.text_frame.text = f"Period: {start_date} to {end_date}"
                for paragraph in shape.text_frame.paragraphs:
                    paragraph.font.name = "Inter"
                    paragraph.font.size = Pt(13)
                    paragraph.font.color.rgb = RGBColor(100, 116, 139)
            elif "lead" in text_lower or "manager" in text_lower:
                shape.text_frame.text = f"Lead: {auditors}"
                for paragraph in shape.text_frame.paragraphs:
                    paragraph.font.name = "Inter"
                    paragraph.font.size = Pt(13)
                    paragraph.font.color.rgb = RGBColor(100, 116, 139)
                    
    # Modify Slide 2 (Scope & Objectives process list)
    slide_2 = prs.slides[1]
    process_list_str = "\n".join([f"- {p['process_name']}" for p in processes_data])
    for shape in slide_2.shapes:
        if shape.has_text_frame:
            if "scope" not in shape.text_frame.text.lower():
                shape.text_frame.text = process_list_str
                for paragraph in shape.text_frame.paragraphs:
                    paragraph.font.name = "Inter"
                    paragraph.font.size = Pt(12)
                    paragraph.font.color.rgb = RGBColor(51, 65, 85)

    # Delete template dummy slides (indices 5, 4, 3, 2)
    if len(prs.slides) >= 7:
        for idx in [5, 4, 3, 2]:
            del prs.slides._sldIdLst[idx]
            
    # Add dynamic slides for each process and its controls
    import math
    import re
    slide_insert_index = 2
    
    # 1. Determine the best slide layout from the template
    # We want a layout named "blank" or similar to avoid ghost placeholders.
    slide_layout = None
    for layout in prs.slide_layouts:
        if layout.name.lower() == 'blank':
            slide_layout = layout
            break
    if slide_layout is None:
        for layout in prs.slide_layouts:
            if 'master' in layout.name.lower():
                slide_layout = layout
                break
    if slide_layout is None:
        slide_layout = prs.slide_layouts[1]
        
    for idx, p_info in enumerate(processes_data):
        nested_controls = p_info["controls"]
        
        # Paginate controls to prevent overflow
        # Card size and text limits
        card_width = prs.slide_width - Inches(1.7)
        card_height = prs.slide_height - Inches(2.4) # Increased height to fit more lines
        
        card_width_in = card_width / Inches(1)
        card_height_in = card_height / Inches(1)
        
        # At font size 14, 1 inch fits about 9.5 characters on average
        text_width_inches = card_width_in - 0.4
        chars_per_line = int(text_width_inches * 9.5)
        
        # Estimate height of controls list in inches dynamically
        def estimate_height_in(controls_list):
            if not controls_list:
                return 0.0
            total_height_pt = 0.0
            for ctrl in controls_list:
                ctrl_str = str(ctrl)
                if ctrl_str.startswith(("•", "-", "*")):
                    ctrl_str = ctrl_str[1:].strip()
                ctrl_lines = [l.strip() for l in re.split(r'[\n\r\x0b]+', ctrl_str) if l.strip()]
                for line in ctrl_lines:
                    lines_needed = max(1, math.ceil(len(line) / chars_per_line))
                    # 14pt font size with ~1.2 line spacing = 16.8pt. Space after is 10pt.
                    total_height_pt += lines_needed * 16.8 + 10.0
            return total_height_pt / 72.0
            
        # Target max text height is card height minus top/bottom padding
        max_text_height_in = card_height_in - 0.3
        
        pages = []
        current_page_controls = []
        
        if not nested_controls:
            pages = [[]]
        else:
            for ctrl in nested_controls:
                temp_page = current_page_controls + [ctrl]
                if estimate_height_in(temp_page) > max_text_height_in and current_page_controls:
                    pages.append(current_page_controls)
                    current_page_controls = [ctrl]
                else:
                    current_page_controls = temp_page
                    
            if current_page_controls:
                pages.append(current_page_controls)
                
        # Create slides for each page
        for page_idx, page_controls in enumerate(pages):
            slide = prs.slides.add_slide(slide_layout)
            
            # Remove any shape that is a placeholder on the slide just in case
            for shape in list(slide.shapes):
                if shape.is_placeholder:
                    sp = shape._element
                    sp.getparent().remove(sp)
            
            # 1. Left-side title accent line
            accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.65), Inches(0.75), Inches(0.06), Inches(0.8))
            accent_bar.fill.solid()
            accent_bar.fill.fore_color.rgb = RGBColor(99, 102, 241) # Indigo accent
            accent_bar.line.color.rgb = RGBColor(99, 102, 241)
            
            # Title text box
            title_box = slide.shapes.add_textbox(Inches(0.9), Inches(0.7), prs.slide_width - Inches(1.8), Inches(0.8))
            tf_title = title_box.text_frame
            tf_title.word_wrap = True
            p_title = tf_title.paragraphs[0]
            
            # Add Part suffix if paginated
            slide_title = p_info['process_name']
            if len(pages) > 1:
                slide_title += f" (Part {page_idx + 1})"
                
            p_title.text = slide_title
            p_title.font.name = "Inter"
            p_title.font.size = Pt(22)
            p_title.font.bold = True
            p_title.font.color.rgb = RGBColor(15, 23, 42)
            
            # 2. Risk & Effort Badges (below title)
            risk_score = p_info.get('risk_score', 50)
            risk_color = RGBColor(239, 68, 68) if risk_score >= 85 else (RGBColor(245, 158, 11) if risk_score >= 70 else RGBColor(16, 185, 129))
            
            # Risk Badge shape
            risk_badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(1.6), Inches(1.8), Inches(0.3))
            risk_badge.fill.solid()
            risk_badge.fill.fore_color.rgb = risk_color
            risk_badge.line.color.rgb = risk_color
            tf_risk = risk_badge.text_frame
            p_risk = tf_risk.paragraphs[0]
            p_risk.text = f"AI RISK SCORE: {risk_score}%"
            p_risk.font.name = "Inter"
            p_risk.font.size = Pt(9)
            p_risk.font.bold = True
            p_risk.font.color.rgb = RGBColor(255, 255, 255)
            tf_risk.margin_left = Inches(0.05)
            tf_risk.margin_right = Inches(0.05)
            tf_risk.margin_top = Inches(0.02)
            tf_risk.margin_bottom = Inches(0.02)
            
            # Effort Badge shape
            hours = p_info.get('man_hours', 40)
            hours_badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.8), Inches(1.6), Inches(1.4), Inches(0.3))
            hours_badge.fill.solid()
            hours_badge.fill.fore_color.rgb = RGBColor(99, 102, 241) # Indigo
            hours_badge.line.color.rgb = RGBColor(99, 102, 241)
            tf_hours = hours_badge.text_frame
            p_hours = tf_hours.paragraphs[0]
            p_hours.text = f"EFFORT: {hours} HRS"
            p_hours.font.name = "Inter"
            p_hours.font.size = Pt(9)
            p_hours.font.bold = True
            p_hours.font.color.rgb = RGBColor(255, 255, 255)
            tf_hours.margin_left = Inches(0.05)
            tf_hours.margin_right = Inches(0.05)
            tf_hours.margin_top = Inches(0.02)
            tf_hours.margin_bottom = Inches(0.02)
            
            # Calculate dynamic card height to fit content nicely
            page_text_height_in = estimate_height_in(page_controls)
            dynamic_card_height_in = page_text_height_in + 0.25 # Top/bottom padding
            dynamic_card_height_in = min(card_height_in, dynamic_card_height_in)
            dynamic_card_height = Inches(dynamic_card_height_in)
            
            # 3. Card Background for Content (dynamic height)
            card_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.85), Inches(2.0), card_width, dynamic_card_height)
            card_bg.fill.solid()
            card_bg.fill.fore_color.rgb = RGBColor(248, 250, 252) # Slate-50 bg
            card_bg.line.color.rgb = RGBColor(226, 232, 240)      # Slate-200 border
            card_bg.line.width = Pt(1.5)
            
            # Content text box (dynamic height matching card)
            content_box = slide.shapes.add_textbox(Inches(1.05), Inches(2.1), card_width - Inches(0.4), dynamic_card_height - Inches(0.2))
            tf_content = content_box.text_frame
            tf_content.word_wrap = True
            tf_content.margin_left = Inches(0)
            tf_content.margin_right = Inches(0)
            tf_content.margin_top = Inches(0)
            tf_content.margin_bottom = Inches(0)
            
            # Bullets for controls
            if not page_controls:
                p_none = tf_content.paragraphs[0]
                p_none.text = "• No key controls assigned to this process area."
                p_none.font.name = "Inter"
                p_none.font.size = Pt(14)
                p_none.font.italic = True
                p_none.font.color.rgb = RGBColor(148, 163, 184)
            else:
                for c_idx, ctrl in enumerate(page_controls):
                    # Clean up bullet character prefix if it already exists in data to avoid duplicate bullets
                    ctrl_str = str(ctrl)
                    if ctrl_str.startswith(("•", "-", "*")):
                        ctrl_str = ctrl_str[1:].strip()
                        
                    ctrl_lines = [l.strip() for l in re.split(r'[\n\r\x0b]+', ctrl_str) if l.strip()]
                    for sp_idx, line in enumerate(ctrl_lines):
                        if len(tf_content.paragraphs) == 1 and tf_content.paragraphs[0].text == "":
                            p_ctrl = tf_content.paragraphs[0]
                        else:
                            p_ctrl = tf_content.add_paragraph()
                            
                        if sp_idx == 0:
                            p_ctrl.text = f"• {line}"
                        else:
                            p_ctrl.text = f"  {line}"
                            
                        p_ctrl.font.name = "Inter"
                        p_ctrl.font.size = Pt(14)  # Increased font size to 14pt!
                        p_ctrl.font.color.rgb = RGBColor(51, 65, 85)
                        p_ctrl.space_after = Pt(10)
                        
            # Move the slide XML element to correct order (index slide_insert_index)
            slide_element = prs.slides._sldIdLst[len(prs.slides) - 1]
            prs.slides._sldIdLst.remove(slide_element)
            prs.slides._sldIdLst.insert(slide_insert_index, slide_element)
            slide_insert_index += 1

    # Logo overlay
    base_dir = os.path.dirname(os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__)))))
    logo_path = os.path.join(base_dir, "ARM", "frontend", "public", "logo2.png")
    if not os.path.exists(logo_path):
        logo_path = r"d:\compliBear\ARM\frontend\public\logo2.png"
        
    for slide in prs.slides:
        if os.path.exists(logo_path):
            try:
                # Top right placement
                logo_width = Inches(1.5)
                # Keep margin of 0.4 inches from right and 0.3 inches from top
                logo_left = prs.slide_width - logo_width - Inches(0.4)
                logo_top = Inches(0.3)
                slide.shapes.add_picture(logo_path, logo_left, logo_top, width=logo_width)
            except Exception as e:
                print(f"Failed to add logo: {e}")

    # Save generated presentation
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    audit_period_safe = f"{start_date}_to_{end_date}".replace(' ', '_')
    out_filename = f"Sentinel_Audit_Plan_{audit_period_safe}.pptx"
    out_path = os.path.join(OUTPUT_DIR, out_filename)
    prs.save(out_path)
    
    return out_path
